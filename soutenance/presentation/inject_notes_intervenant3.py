# -*- coding: utf-8 -*-
"""Injecte le discours de l'Intervenant 3 dans les NOTES du presentateur
des diapos 12 a 18 du PPTX de soutenance.

- Mapping verifie : index python-pptx 12..18 == pied de page 12..18.
- Idempotent : remplace toujours la note par la version ci-dessous.
- Sauvegarde .BAK horodatee avant ecriture.
"""
import os
import shutil
import datetime
from pptx import Presentation

HERE = os.path.dirname(os.path.abspath(__file__))
PPTX = os.path.join(HERE, "CuraMedical-Présentation.pptx")

# index slide -> (entete, timing, texte du discours)
NOTES = {
    12: ("Ordonnances & documents", "~ 0:47",
         "Poursuivons. Pour nous, une consultation ne se termine pas quand le médecin "
         "referme le dossier : elle se termine quand le patient repart avec un vrai "
         "document en main. Et c'est exactement ce que CuraMedical produit, tout seul. "
         "Le médecin saisit son ordonnance de façon structurée : pour chaque "
         "médicament, la dose et la durée. À partir de cette saisie, l'ordonnance "
         "est générée automatiquement en PDF, avec le cachet et la signature du "
         "médecin déjà intégrés — le document est prêt à imprimer ou à envoyer. "
         "Les comptes-rendus de consultation suivent exactement la même logique et "
         "sont eux aussi exportés en PDF. On passe donc de la saisie à un document "
         "propre et professionnel en un seul clic, sans jamais rien retaper, et surtout "
         "sans risque d'erreur de recopie."),

    13: ("Téléconsultation", "~ 0:37",
         "Mais tous les patients ne peuvent pas se déplacer jusqu'au cabinet — par "
         "distance, ou par difficulté à se déplacer. Alors, plutôt que de faire "
         "venir le patient, CuraMedical amène le cabinet jusqu'à lui, grâce à la "
         "téléconsultation. La visioconférence est intégrée directement dans la "
         "plateforme via Jitsi : il n'y a rien à installer. Pour chaque rendez-vous "
         "à distance, une salle privée est créée automatiquement, et pendant "
         "l'appel, le médecin prend ses notes en direct. L'accès aux soins n'est "
         "plus une question de kilomètres."),

    14: ("Chatbot IA", "~ 0:41",
         "Maintenant, posons-nous une question toute simple : pour savoir combien de "
         "patients ont un rendez-vous aujourd'hui, faut-il vraiment cliquer à travers "
         "cinq écrans ? Nous avons décidé que non. Nous avons donc ajouté un chatbot "
         "qui permet d'interroger le cabinet en langage naturel. On tape sa question en "
         "français courant — « combien de patients avons-nous ? », « quels sont "
         "les rendez-vous d'aujourd'hui ? » — et la réponse arrive immédiatement, "
         "calculée sur les vraies données de l'application. Ce chatbot s'appuie sur "
         "un grand modèle de langage, via Groq. C'est donc une deuxième utilisation de "
         "l'IA dans le projet : non plus pour le diagnostic, mais pour fluidifier la "
         "gestion de tous les jours."),

    15: ("Notifications WhatsApp", "~ 0:41",
         "Pour la communication avec le patient, nous avons fait un constat tout simple : "
         "un beau site web, ce n'est utile que si le patient sait s'en servir. Or tout "
         "le monde n'est pas à l'aise pour se connecter à un site et naviguer dans une "
         "interface — je pense en particulier aux personnes âgées. Nous avons donc "
         "décidé d'aller chercher le patient là où il se trouve déjà, sur "
         "l'application que tout le monde sait utiliser : WhatsApp. Concrètement, le "
         "patient reçoit ses rappels de rendez-vous automatiquement, ce qui évite "
         "énormément d'oublis, et il peut même recevoir son ordonnance directement "
         "sur son téléphone. Le tout est orchestré par un scénario automatique "
         "construit avec n8n et relié à Twilio. Le résultat est double : moins de "
         "rendez-vous manqués pour le cabinet, et un patient mieux suivi, plus "
         "rassuré — sans qu'il ait quoi que ce soit à apprendre."),

    16: ("Espace patient", "~ 0:39",
         "Et le patient, justement, n'est plus un simple spectateur de son suivi : il en "
         "devient acteur. Depuis son espace personnel, il peut s'inscrire en ligne, "
         "consulter ses prochains rendez-vous, retrouver ses consultations et ses "
         "ordonnances, et gérer ses informations. Il dispose ainsi, à tout moment et où "
         "qu'il soit, d'une vue claire sur son parcours de soins. Et là, la boucle "
         "est complète : de la secrétaire au médecin, et du médecin jusqu'au "
         "patient."),

    17: ("Difficultés techniques rencontrées", "~ 0:50",
         "Avant de conclure, un mot sur les difficultés rencontrées — parce qu'un "
         "projet, ce sont aussi des obstacles qu'on apprend à franchir. Côté IA, le "
         "plus délicat a été d'aligner un jeu de données en anglais avec une saisie "
         "en français. Côté sécurité, il fallait cloisonner les données par rôle, "
         "sans pour autant rendre l'application pénible à utiliser. Pour les "
         "documents, produire des PDF à la fois lisibles et réutilisables nous a "
         "demandé plusieurs essais. L'automatisation des envois nous a poussés à "
         "découpler les tâches avec Celery, n8n et Twilio, pour ne jamais bloquer "
         "l'application. Enfin, faire tourner ensemble plusieurs services "
         "conteneurisés, et garder une interface claire malgré quatre profils "
         "différents, ont été de vrais défis. Chacun a été analysé, puis résolu "
         "— et c'est très honnêtement ce qui nous a le plus fait progresser."),

    18: ("Bilan & perspectives", "~ 0:56",
         "Pour conclure. Nous sommes partis d'un problème concret — un cabinet "
         "désorganisé, sans aucune aide à la décision — et nous avons abouti à "
         "une plateforme complète et réellement fonctionnelle : quatre rôles bien "
         "séparés et sécurisés, une intelligence artificielle véritablement "
         "opérationnelle, des documents générés automatiquement, et de vraies "
         "connexions avec des outils externes. Ce projet nous a fait vivre toutes les "
         "étapes du métier d'ingénieur, de l'analyse du besoin jusqu'à la mise en "
         "service. Et nous ne comptons pas nous arrêter là : nous aimerions déployer "
         "la plateforme en production, continuer à améliorer le modèle d'IA, créer "
         "une application mobile pour les patients, et la relier à d'autres systèmes "
         "de santé. Merci beaucoup pour votre attention — nous sommes maintenant "
         "prêts à répondre à vos questions."),
}


def main():
    ts = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
    bak = PPTX.replace(".pptx", f".BAK_notes_{ts}.pptx")
    shutil.copy2(PPTX, bak)
    print("Sauvegarde :", os.path.basename(bak))

    prs = Presentation(PPTX)
    for idx, (entete, timing, texte) in NOTES.items():
        slide = prs.slides[idx]
        note = (f"INTERVENANT 3 — Slide {idx} / 18 — {entete}  ({timing})\n\n"
                f"{texte}")
        slide.notes_slide.notes_text_frame.text = note
        print(f"  notes injectees -> idx {idx} ({entete})")
    prs.save(PPTX)
    print("OK :", os.path.basename(PPTX))


if __name__ == "__main__":
    main()
