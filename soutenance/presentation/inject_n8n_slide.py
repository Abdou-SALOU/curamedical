# -*- coding: utf-8 -*-
"""Injecte l'image du workflow n8n (capture de l'editeur) dans la diapo
« Automatisation · n8n » du livrable CuraMedical-Presentation.pptx, EN PLACE,
en preservant les retouches manuelles (pas de regeneration).

- repere la diapo n8n via l'eyebrow « AUTOMATISATION · N8N » ;
- reduit la capture WhatsApp (portrait) en encart a droite ;
- ajoute le workflow (paysage) en visuel central, sur une carte + legende ;
- idempotent (ne reinjecte pas) ; sauvegarde .bak avant ecriture.

Usage : python inject_n8n_slide.py
"""
import os
import shutil

from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE

HERE = os.path.dirname(os.path.abspath(__file__))
DECK = os.path.join(HERE, "CuraMedical-Présentation.pptx")
WF = os.path.join(HERE, "images", "17-workflow-n8n.png")

WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0xE2, 0xE8, 0xF0)
GREY = RGBColor(0x64, 0x74, 0x8B)
FONT = "Segoe UI"
CAP = "Éditeur n8n — workflow d'automatisation réel"


def is_n8n_slide(s):
    for sh in s.shapes:
        if sh.has_text_frame and "automatisation · n8n" in sh.text_frame.text.strip().lower():
            return True
    return False


def already_done(s):
    for sh in s.shapes:
        if sh.has_text_frame and CAP[:18].lower() in sh.text_frame.text.strip().lower():
            return True
    return False


def find_phone(s):
    """La capture WhatsApp = grande image portrait (hauteur > largeur, h > 4 po)."""
    best = None
    for sh in s.shapes:
        if sh.shape_type == 13:  # PICTURE
            if sh.height > sh.width and sh.height > Inches(4):
                if best is None or sh.height > best.height:
                    best = sh
    return best


def main():
    assert os.path.exists(DECK), f"deck introuvable : {DECK}"
    assert os.path.exists(WF), f"image introuvable : {WF}"

    prs = Presentation(DECK)
    target = None
    for i, s in enumerate(prs.slides):
        if is_n8n_slide(s):
            target = (i, s)
            break
    if target is None:
        print("DIAPO n8n NON TROUVEE")
        return
    idx, s = target
    if already_done(s):
        print(f"DEJA FAIT (diapo {idx})")
        return

    # 1) reduire + deplacer la capture WhatsApp en encart a droite
    phone = find_phone(s)
    if phone is not None:
        new_w = Inches(1.78)
        ratio = phone.width / phone.height
        phone.width = new_w
        phone.height = Inches(1.78 / ratio)
        phone.left = Inches(10.85)
        phone.top = Inches(2.7)

    # 2) workflow paysage au centre, sur une carte blanche + legende
    wf_iw, wf_ih = Image.open(WF).size
    ar = wf_iw / wf_ih
    wf_w = 4.4
    wf_h = wf_w / ar
    wf_x, wf_y = 6.3, 2.95
    pad = 0.07
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(wf_x - pad), Inches(wf_y - pad),
                              Inches(wf_w + 2 * pad), Inches(wf_h + 2 * pad))
    card.fill.solid(); card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = LINE; card.line.width = Pt(1)
    card.shadow.inherit = False

    pic = s.shapes.add_picture(WF, Inches(wf_x), Inches(wf_y), Inches(wf_w), Inches(wf_h))
    pic.line.color.rgb = LINE; pic.line.width = Pt(0.75)
    pic.shadow.inherit = False

    cap_tb = s.shapes.add_textbox(Inches(wf_x - pad), Inches(wf_y + wf_h + 0.10),
                                  Inches(wf_w + 2 * pad), Inches(0.3))
    tf = cap_tb.text_frame
    tf.word_wrap = True; tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = CAP
    r.font.size = Pt(10.5); r.font.italic = True; r.font.name = FONT; r.font.color.rgb = GREY

    bak = DECK.replace(".pptx", ".bak_n8n.pptx")
    shutil.copyfile(DECK, bak)
    prs.save(DECK)
    print(f"OK -> diapo {idx} : workflow ajoute (carte + legende), capture WhatsApp reduite.")
    print(f"Sauvegarde : {os.path.basename(bak)}")


if __name__ == "__main__":
    main()
