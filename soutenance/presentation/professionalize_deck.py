# -*- coding: utf-8 -*-
"""Améliorations ciblées de la présentation CuraMedical.

Le script préserve le fichier original via une sauvegarde, corrige quelques
formulations trop absolues et ajoute deux slides de synthèse professionnelle :
parcours utilisateur et défis techniques.
"""
from __future__ import annotations

import shutil
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
DECK = ROOT / "CuraMedical-Présentation.pptx"
BACKUP = ROOT / "CuraMedical-Présentation.backup-professionnelle.pptx"
OUTPUT = ROOT / "CuraMedical-Présentation-professionnelle.pptx"

BG = RGBColor(0xF7, 0xF8, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x2A, 0x9B, 0x69)
DARK = RGBColor(0x0F, 0x17, 0x2A)
GREY = RGBColor(0x64, 0x74, 0x8B)
LINE = RGBColor(0xE2, 0xE8, 0xF0)
VIOLET = RGBColor(0x89, 0x00, 0xFF)
VIOLBG = RGBColor(0xF3, 0xE8, 0xFF)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
AMBER_BG = RGBColor(0xFF, 0xF7, 0xED)
FONT = "Segoe UI"


def set_bg(slide, rgb=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def rect(slide, x, y, w, h, fill=None, line=None, rounded=False, lw=Pt(1)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        x,
        y,
        w,
        h,
    )
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = lw
    shape.shadow.inherit = False
    return shape


def tbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    return tf


def par(tf, runs, size=14, color=DARK, bold=False, first=False, align=PP_ALIGN.LEFT, sa=None, ls=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    if sa is not None:
        p.space_after = Pt(sa)
    if ls is not None:
        p.line_spacing = ls
    if isinstance(runs, str):
        runs = [(runs, {})]
    for text, opts in runs:
        run = p.add_run()
        run.text = text
        font = run.font
        font.name = FONT
        font.size = Pt(opts.get("size", size))
        font.bold = opts.get("bold", bold)
        font.color.rgb = opts.get("color", color)
    return p


def card(slide, x, y, w, h, fill=WHITE, line=LINE):
    return rect(slide, Inches(x), Inches(y), Inches(w), Inches(h), fill=fill, line=line, rounded=True)


def header(slide, eyebrow, title, subtitle=None, accent=GREEN):
    tf = tbox(slide, Inches(0.62), Inches(0.52), Inches(12), Inches(0.3))
    par(tf, [(eyebrow.upper(), {"size": 12.5, "bold": True, "color": accent})], 12.5, accent, True, first=True)
    tt = tbox(slide, Inches(0.6), Inches(0.84), Inches(12.1), Inches(0.95))
    par(tt, [(title, {"size": 38, "bold": True, "color": DARK})], 38, DARK, True, first=True)
    if subtitle:
        st = tbox(slide, Inches(0.62), Inches(1.66), Inches(11.8), Inches(0.6))
        par(st, [(subtitle, {"size": 17, "color": GREY})], 17, GREY, first=True)


def footer(slide, n, total):
    rect(slide, Inches(0.6), Inches(7.02), Inches(12.13), Pt(1), fill=LINE)
    icon = ROOT / "images" / "00-icone-curamedical.png"
    if icon.exists():
        slide.shapes.add_picture(str(icon), Inches(0.6), Inches(7.12), height=Inches(0.24))
    left = tbox(slide, Inches(0.95), Inches(7.13), Inches(3.8), Inches(0.3))
    par(left, [("CuraMedical", {"size": 10.5, "bold": True, "color": GREEN})], 10.5, GREEN, True, first=True)
    right = tbox(slide, Inches(8.0), Inches(7.13), Inches(4.73), Inches(0.3))
    par(right, [(f"{n:02d} / {total}", {"size": 10.5, "color": GREY})], 10.5, GREY, align=PP_ALIGN.RIGHT, first=True)


def replace_text(prs, replacements):
    count = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text = run.text
                    for old, new in replacements.items():
                        if old in text:
                            run.text = text.replace(old, new)
                            count += 1
    return count


def rewrite_bullet_shape(shape, items, mark_color=GREEN):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = MSO_ANCHOR.TOP
    first = True
    for head, body in items:
        par(
            tf,
            [
                ("›  ", {"size": 16, "bold": True, "color": mark_color}),
                (head + " ", {"size": 16, "bold": True, "color": DARK}),
                (body, {"size": 16, "color": GREY}),
            ],
            16,
            DARK,
            first=first,
            sa=11,
            ls=1.12,
        )
        first = False


def rewrite_sensitive_slides(prs):
    rewrites = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = " ".join(shape.text_frame.text.split())
            if (
                "Authentification par JWT" in text
                and "Cloisonnement strict" in text
                and "Corbeille" in text
            ):
                rewrite_bullet_shape(
                    shape,
                    [
                        ("Authentification par JWT", "(SimpleJWT)."),
                        ("Permissions par rôle", "au niveau de l'API (DRF)."),
                        ("Contrôle d'accès strict", "selon les rôles et les périmètres fonctionnels."),
                        ("Corbeille et journalisation", "des principales entités sensibles."),
                    ],
                )
                rewrites += 1
            elif "accès réservé aux professionnels" in text and "journalisées" in text:
                shape.text_frame.text = (
                    "Porte d’entrée sécurisée — accès réservé aux professionnels, "
                    "les actions sensibles sont journalisées."
                )
                rewrites += 1
            elif "Questions en français courant" in text and "Propulsé par un LLM" in text:
                rewrite_bullet_shape(
                    shape,
                    [
                        ("Questions en français courant", ": statistiques, patients, rendez-vous."),
                        ("Réponses instantanées", ": nombre de patients, RDV du jour..."),
                        ("Propulsé par Groq", "(Llama 3, avec repli local)."),
                    ],
                )
                rewrites += 1
    return rewrites


def add_journey_slide(prs, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    header(
        slide,
        "Démonstration · Parcours métier",
        "Le parcours complet de soins",
        "Une lecture simple du projet : de la demande patient jusqu'au document transmis.",
    )

    steps = [
        ("01", "Inscription", "Le patient crée son compte ou est ajouté par le secrétariat."),
        ("02", "Rendez-vous", "Le créneau est planifié, validé et notifié."),
        ("03", "Consultation", "Le médecin saisit les symptômes et l'examen clinique."),
        ("04", "Suggestion IA", "Le modèle propose des hypothèses, le médecin décide."),
        ("05", "Documents", "Ordonnance et compte-rendu sont générés en PDF."),
        ("06", "Notification", "Le patient reçoit les informations par e-mail/WhatsApp."),
    ]
    x0, y0, w, h, gap = 0.6, 2.55, 1.85, 2.9, 0.22
    for i, (num, title, body) in enumerate(steps):
        x = x0 + i * (w + gap)
        card(slide, x, y0, w, h)
        tf = tbox(slide, Inches(x + 0.18), Inches(y0 + 0.22), Inches(w - 0.36), Inches(h - 0.44))
        par(tf, [(num, {"size": 24, "bold": True, "color": GREEN})], 24, GREEN, True, first=True, sa=10)
        par(tf, [(title, {"size": 15, "bold": True, "color": DARK})], 15, DARK, True, sa=7)
        par(tf, [(body, {"size": 11.2, "color": GREY})], 11.2, GREY, ls=1.05)
        if i < len(steps) - 1:
            arrow = tbox(slide, Inches(x + w - 0.02), Inches(y0 + 1.18), Inches(gap + 0.12), Inches(0.35), MSO_ANCHOR.MIDDLE)
            par(arrow, [("→", {"size": 17, "bold": True, "color": GREEN})], 17, GREEN, True, first=True, align=PP_ALIGN.CENTER)

    band = card(slide, 1.25, 5.95, 10.85, 0.62, fill=VIOLBG, line=VIOLBG)
    tf = tbox(slide, Inches(1.55), Inches(6.07), Inches(10.25), Inches(0.36), MSO_ANCHOR.MIDDLE)
    par(
        tf,
        [
            ("Message à retenir : ", {"size": 13.5, "bold": True, "color": VIOLET}),
            ("CuraMedical couvre un flux métier complet, pas seulement une collection d'écrans.", {"size": 13.5, "color": DARK}),
        ],
        13.5,
        DARK,
        first=True,
    )
    footer(slide, total - 1, total)
    return slide


def add_challenges_slide(prs, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    header(
        slide,
        "Ingénierie · Difficultés",
        "Défis techniques maîtrisés",
        "Les choix techniques importants qui montrent la valeur d'ingénierie du projet.",
        accent=AMBER,
    )

    items = [
        ("IA médicale", "Aligner un dataset anglophone avec une saisie utilisateur en français."),
        ("Sécurité", "Cloisonner les données selon les rôles sans bloquer les parcours métier."),
        ("PDF médicaux", "Produire des documents lisibles, signés et réutilisables côté patient."),
        ("Automatisation", "Découpler les envois via Celery, n8n, Twilio et e-mail."),
        ("Déploiement", "Faire fonctionner plusieurs services conteneurisés ensemble."),
        ("Expérience", "Garder une interface claire malgré quatre profils utilisateurs."),
    ]
    for i, (title, body) in enumerate(items):
        row, col = divmod(i, 2)
        x = 0.75 + col * 6.0
        y = 2.35 + row * 1.35
        fill = AMBER_BG if i in (0, 3) else WHITE
        card(slide, x, y, 5.5, 1.08, fill=fill, line=LINE)
        tf = tbox(slide, Inches(x + 0.26), Inches(y + 0.18), Inches(5.0), Inches(0.8))
        par(tf, [(title, {"size": 14.5, "bold": True, "color": DARK})], 14.5, DARK, True, first=True, sa=3)
        par(tf, [(body, {"size": 11.5, "color": GREY})], 11.5, GREY, ls=1.03)

    tf = tbox(slide, Inches(0.8), Inches(6.35), Inches(11.7), Inches(0.42), MSO_ANCHOR.MIDDLE)
    par(
        tf,
        [("Cette slide aide le jury à voir le travail réel derrière l'interface.", {"size": 13, "bold": True, "color": AMBER})],
        13,
        AMBER,
        True,
        first=True,
        align=PP_ALIGN.CENTER,
    )
    footer(slide, total, total)
    return slide


def move_slide(prs, old_index, new_index):
    slide_id_list = prs.slides._sldIdLst
    slides = list(slide_id_list)
    slide = slides[old_index]
    slide_id_list.remove(slide)
    slide_id_list.insert(new_index, slide)


def renumber_footers(prs, total):
    # Update existing page numbers in place. Add a footer only to new slides.
    for index, slide in enumerate(list(prs.slides)[1:], start=1):
        found = False
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if "/" in text and len(text) <= 8:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if "/" in run.text:
                            run.text = f"{index:02d} / {total}"
                            found = True
                if not found:
                    shape.text_frame.text = f"{index:02d} / {total}"
                    found = True
        if not found:
            footer(slide, index, total)


def main():
    if not DECK.exists():
        raise FileNotFoundError(DECK)
    shutil.copyfile(DECK, BACKUP)
    prs = Presentation(DECK)

    replacements = {
        "Une architecture moderne, en microservices, entièrement conteneurisée.": "Une architecture hybride, conteneurisée : backend Django modulaire, microservice IA et automatisations.",
        "Une architecture moderne, en microservices, entiÃ¨rement conteneurisÃ©e.": "Une architecture hybride, conteneurisée : backend Django modulaire, microservice IA et automatisations.",
        "Cloisonnement strict\ndes données ; tout accès non autorisé est bloqué.": "Contrôle d'accès strict\nselon les rôles et les périmètres fonctionnels.",
        "Cloisonnement strict\ndes donnÃ©es ; tout accÃ¨s non autorisÃ© est bloquÃ©.": "Contrôle d'accès strict\nselon les rôles et les périmètres fonctionnels.",
        "Corbeille et traçabilité\ndes actions.": "Corbeille et journalisation\ndes principales entités sensibles.",
        "Corbeille et traÃ§abilitÃ©\ndes actions.": "Corbeille et journalisation\ndes principales entités sensibles.",
        "de précision sur le jeu de test": "de précision sur 1 500 cas de test",
        "de prÃ©cision sur le jeu de test": "de précision sur 1 500 cas de test",
        "Propulsé par un LLM\n(Groq).": "Propulsé par Groq\n(Llama 3 + repli local).",
        "PropulsÃ© par un LLM\n(Groq).": "Propulsé par Groq\n(Llama 3 + repli local).",
    }
    changed = replace_text(prs, replacements)
    rewrites = rewrite_sensitive_slides(prs)

    # Total hors couverture : les 18 slides existantes après couverture + 2 nouvelles.
    total = len(prs.slides) - 1 + 2
    add_journey_slide(prs, total)
    move_slide(prs, len(prs.slides) - 1, 10)  # avant la partie IA.
    add_challenges_slide(prs, total)
    move_slide(prs, len(prs.slides) - 1, len(prs.slides) - 3)  # avant conclusion/merci.

    renumber_footers(prs, total)
    prs.save(OUTPUT)
    print(f"OK: {OUTPUT}")
    print(f"Sauvegarde: {BACKUP}")
    print(f"Textes corrigés: {changed}")
    print(f"Slides réécrites: {rewrites}")
    print(f"Slides finales: {len(prs.slides)}")


if __name__ == "__main__":
    main()
