# -*- coding: utf-8 -*-
"""Ajoute la diapo finale "Merci beaucoup pour votre attention" :
   image plein ecran (16:9) + bandeau translucide bas + texte. Diapo NON numerotee."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls, qn

DOCX = "CuraMedical-Présentation.pptx"
IMG  = "images/merci-illustration.png"

# palette (charte CuraMedical)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREEN  = RGBColor(0x2A, 0x9B, 0x69)
GREENL = RGBColor(0x4E, 0xB5, 0x83)
DARK   = RGBColor(0x0F, 0x17, 0x2A)
LGREY  = RGBColor(0xCB, 0xD5, 0xE1)
FONT   = "Segoe UI"

prs = Presentation(DOCX)
SW, SH = prs.slide_width, prs.slide_height   # 13.333 x 7.5 in EMU

# --- etat avant ---
print("AVANT : %d diapos" % len(prs.slides._sldIdLst))
for i, s in enumerate(prs.slides, 1):
    t = ""
    for sh in s.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip():
            t = sh.text_frame.text.strip().split("\n")[0]; break
    print("  %2d. %s" % (i, t[:48]))

BLANK = prs.slide_layouts[6]
s = prs.slides.add_slide(BLANK)

def set_alpha(shape, opacity_pct):
    """opacity_pct : 0 (transparent) -> 100 (opaque)"""
    srgb = shape._element.spPr.find(qn('a:solidFill')).find(qn('a:srgbClr'))
    srgb.append(parse_xml('<a:alpha %s val="%d"/>' % (nsdecls('a'), int(opacity_pct*1000))))

def tbox(x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.auto_size = MSO_AUTO_SIZE.NONE; tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    return tf

def par(tf, runs, size, color, bold=False, align=PP_ALIGN.CENTER, first=False, sa=None, ls=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    if sa is not None: p.space_after = Pt(sa)
    if ls is not None: p.line_spacing = ls
    for t, o in runs:
        r = p.add_run(); r.text = t; f = r.font
        f.size = Pt(o.get('size', size)); f.bold = o.get('bold', bold)
        f.name = FONT; f.color.rgb = o.get('color', color)
    return p

# --- 1) image plein ecran (16:9 = format de la diapo, aucune deformation) ---
s.shapes.add_picture(IMG, 0, 0, width=SW, height=SH)

# --- 2) bandeau translucide en bas pour la lisibilite ---
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(5.18), SW, SH - Inches(5.18))
band.fill.solid(); band.fill.fore_color.rgb = DARK
band.line.fill.background(); band.shadow.inherit = False
set_alpha(band, 82)   # 82% opaque

# fin liseret vert haut du bandeau
acc = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(5.18), SW, Pt(3))
acc.fill.solid(); acc.fill.fore_color.rgb = GREEN
acc.line.fill.background(); acc.shadow.inherit = False

# --- 3) texte centre dans le bandeau ---
tf = tbox(Inches(0.6), Inches(5.18), Inches(12.13), SH - Inches(5.18), anchor=MSO_ANCHOR.MIDDLE)
par(tf, [("CURAMEDICAL  ·  SANTÉ NUMÉRIQUE PAR L'IA",
          {'size': 12.5, 'bold': True, 'color': GREENL})], 12.5, GREENL, first=True, sa=10)
par(tf, [("Merci beaucoup ", {'size': 36, 'bold': True, 'color': WHITE}),
         ("pour votre attention", {'size': 36, 'bold': True, 'color': GREENL})],
    36, WHITE, sa=12)
par(tf, [("Groupe 2  ·  Abdou SALOU ABDOU  ·  Kamara MACIRE  ·  Nouridine SAWADOGO",
          {'size': 13, 'color': LGREY})], 13, LGREY)

prs.save(DOCX)
print("\nAPRES : %d diapos -> diapo 'Merci' ajoutee en derniere position." % len(prs.slides._sldIdLst))
