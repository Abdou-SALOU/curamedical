# -*- coding: utf-8 -*-
"""Supprime la diapo "Demonstration" et remet les pieds de page en NN / 17."""
import re
from pptx import Presentation
from pptx.oxml.ns import qn

DOCX = "CuraMedical-Présentation.pptx"
TOTAL = 17

prs = Presentation(DOCX)

# reperer la diapo demo
demo = None
for s in prs.slides:
    for sh in s.shapes:
        if sh.has_text_frame:
            t = sh.text_frame.text
            if "Place à la démonstration" in t or ("DÉMONSTRATION" in t.upper() and "LIVE" in t.upper()):
                demo = s; break
    if demo: break
assert demo is not None, "diapo demo introuvable"

# rId de la diapo dans la presentation
demo_part = demo.part
rId = None
for rid, rel in prs.part.rels.items():
    if rel.target_part is demo_part:
        rId = rid; break
assert rId, "relation de la diapo introuvable"

# retirer du sldIdLst + relation
sldIdLst = prs.slides._sldIdLst
for sldId in list(sldIdLst):
    if sldId.get(qn('r:id')) == rId:
        sldIdLst.remove(sldId); break
try:
    prs.part.rels._rels.pop(rId, None)
except Exception:
    pass
print(f"Diapo demo supprimee (rId={rId}).")

# renumeroter les pieds de page NN / 17
pat = re.compile(r'^\s*\d{1,2}\s*/\s*\d{1,2}\s*$')
n = 0
for i, s in enumerate(prs.slides):
    if i == 0:
        continue
    for sh in s.shapes:
        if not sh.has_text_frame:
            continue
        for para in sh.text_frame.paragraphs:
            joined = ''.join(r.text for r in para.runs)
            if pat.match(joined):
                para.runs[0].text = "%02d / %d" % (i, TOTAL)
                for r in para.runs[1:]:
                    r.text = ""
                n += 1
print(f"Pieds de page renumerotes : {n}")

prs.save(DOCX)
print("SAVED. Total diapos:", len(prs.slides._sldIdLst))
