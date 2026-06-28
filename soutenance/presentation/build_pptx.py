# -*- coding: utf-8 -*-
"""Génère CuraMedical-CC.pptx : 17 slides, style éditorial moderne, sans bug.
Tout le texte a word_wrap=True et des tailles maîtrisées -> aucun texte coupé."""
import os
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE

IMG = "images"
def P(name): return os.path.join(IMG, name)

# ---- Palette ----
BG     = RGBColor(0xF7, 0xF8, 0xFA)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREEN  = RGBColor(0x2A, 0x9B, 0x69)
GREENL = RGBColor(0x4E, 0xB5, 0x83)
DARK   = RGBColor(0x0F, 0x17, 0x2A)
GREY   = RGBColor(0x64, 0x74, 0x8B)
LINE   = RGBColor(0xE2, 0xE8, 0xF0)
VIOLET = RGBColor(0x89, 0x00, 0xFF)
VIOLBG = RGBColor(0xF3, 0xE8, 0xFF)
DARKBG = RGBColor(0x0F, 0x14, 0x19)
CARDDK = RGBColor(0x1A, 0x22, 0x30)
LGREY  = RGBColor(0x94, 0xA3, 0xB8)
FONT   = "Segoe UI"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
TOTAL = 17

# ---------- helpers ----------
def slide():
    s = prs.slides.add_slide(BLANK)
    return s

def set_bg(s, rgb):
    f = s.background.fill; f.solid(); f.fore_color.rgb = rgb

def rect(s, x, y, w, h, fill=None, line=None, lw=Pt(1), rounded=False):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb = line; shp.line.width = lw
    shp.shadow.inherit = False
    return shp

def hline(s, x, y, w, color=LINE, weight=1.3):
    return rect(s, x, y, w, Pt(weight), fill=color)

def tbox(s, x, y, w, h, anchor=MSO_ANCHOR.TOP, wrap=True):
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = wrap; tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    return tf

def par(tf, runs, size, color, bold=False, align=PP_ALIGN.LEFT,
        first=False, sa=None, ls=None, font=FONT):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    if sa is not None: p.space_after = Pt(sa)
    if ls is not None: p.line_spacing = ls
    if isinstance(runs, str): runs = [(runs, {})]
    for t, o in runs:
        r = p.add_run(); r.text = t; f = r.font
        f.size = Pt(o.get('size', size)); f.bold = o.get('bold', bold)
        f.name = o.get('font', font); f.color.rgb = o.get('color', color)
    return p

def img_fit(s, path, x, y, w, h, card=True):
    if not os.path.exists(path): return None
    iw, ih = Image.open(path).size
    ar = iw / ih; box = w / h
    if ar > box: nw, nh = w, int(w / ar)
    else: nh, nw = h, int(h * ar)
    nx = x + (w - nw) // 2; ny = y + (h - nh) // 2
    if card:
        pad = Inches(0.07)
        rect(s, nx - pad, ny - pad, nw + 2 * pad, nh + 2 * pad,
             fill=WHITE, line=LINE, lw=Pt(1), rounded=True)
    pic = s.shapes.add_picture(path, nx, ny, nw, nh)
    pic.line.color.rgb = LINE; pic.line.width = Pt(0.75)
    pic.shadow.inherit = False
    return pic

def eyebrow(s, txt, color=GREEN, y=0.52):
    tf = tbox(s, Inches(0.62), Inches(y), Inches(12), Inches(0.3))
    par(tf, [(txt.upper(), {'size': 12.5, 'bold': True, 'color': color})],
        12.5, color, True, first=True)

def title(s, runs, color=DARK, y=0.84):
    tf = tbox(s, Inches(0.6), Inches(y), Inches(12.1), Inches(0.95))
    par(tf, runs, 38, color, True, first=True)

def keyline(s, txt, color=GREY, y=1.66):
    tf = tbox(s, Inches(0.62), Inches(y), Inches(11.8), Inches(0.6))
    par(tf, [(txt, {'size': 17, 'color': color})], 17, color, first=True)

def bullets(s, items, x=0.62, y=2.45, w=5.55, h=4.1, size=16,
            mark=GREEN, head=DARK, body=GREY, sa=11):
    tf = tbox(s, Inches(x), Inches(y), Inches(w), Inches(h))
    first = True
    for it in items:
        runs = [("›  ", {'color': mark, 'bold': True, 'size': size})]
        if isinstance(it, tuple):
            b, n = it
            if b: runs.append((b + " ", {'bold': True, 'color': head, 'size': size}))
            if n: runs.append((n, {'color': body, 'size': size}))
        else:
            runs.append((it, {'color': head, 'size': size}))
        par(tf, runs, size, head, first=first, sa=sa, ls=1.12); first = False

def footer(s, n, dark=False):
    lc = RGBColor(0x24, 0x30, 0x40) if dark else LINE
    tc = LGREY if dark else GREY
    hline(s, Inches(0.6), Inches(7.02), Inches(12.13), color=lc, weight=1)
    if os.path.exists(P("00-icone-curamedical.png")):
        s.shapes.add_picture(P("00-icone-curamedical.png"),
                             Inches(0.6), Inches(7.12), height=Inches(0.24))
    tf = tbox(s, Inches(0.95), Inches(7.13), Inches(5), Inches(0.3))
    par(tf, [("CuraMedical", {'size': 10.5, 'bold': True, 'color': GREEN})],
        10.5, GREEN, True, first=True)
    tf2 = tbox(s, Inches(8.0), Inches(7.13), Inches(4.73), Inches(0.3))
    par(tf2, [("%02d / %d" % (n, TOTAL), {'size': 10.5, 'color': tc})],
        10.5, tc, align=PP_ALIGN.RIGHT, first=True)

def card(s, x, y, w, h, fill=WHITE, line=LINE):
    return rect(s, Inches(x), Inches(y), Inches(w), Inches(h),
                fill=fill, line=line, lw=Pt(1), rounded=True)

# ---------- standard image + bullets slide ----------
def std(n, eb, t_runs, key, items, images=(), ib=(6.55, 2.35, 6.18, 4.35),
        two='stack'):
    s = slide(); set_bg(s, BG)
    eyebrow(s, eb); title(s, t_runs); keyline(s, key); bullets(s, items)
    x, y, w, h = ib
    if len(images) == 1:
        img_fit(s, images[0], Inches(x), Inches(y), Inches(w), Inches(h))
    elif len(images) == 2 and two == 'stack':
        g = 0.22; eh = (h - g) / 2
        img_fit(s, images[0], Inches(x), Inches(y), Inches(w), Inches(eh))
        img_fit(s, images[1], Inches(x), Inches(y + eh + g), Inches(w), Inches(eh))
    footer(s, n)
    return s

# =================== SLIDE 1 — COVER ===================
s = slide(); set_bg(s, BG)
if os.path.exists(P("00-logo-curamedical.png")):
    s.shapes.add_picture(P("00-logo-curamedical.png"), Inches(0.6), Inches(0.55),
                         height=Inches(0.62))
mt = tbox(s, Inches(7.5), Inches(0.55), Inches(5.23), Inches(0.7))
par(mt, [("PROJET EN SYSTÈMES INFORMATIQUES — 2026", {'size': 11, 'color': GREY})],
    11, GREY, align=PP_ALIGN.RIGHT, first=True)
par(mt, [("ISGA · 2CI-ISI · GROUPE 2", {'size': 11, 'color': LGREY})],
    11, LGREY, align=PP_ALIGN.RIGHT)
eyebrow(s, "Santé numérique · Aide au diagnostic par IA", y=2.15)
tt = tbox(s, Inches(0.58), Inches(2.5), Inches(12), Inches(1.5))
par(tt, [("Cura", {'size': 72, 'bold': True, 'color': DARK}),
         ("Medical", {'size': 72, 'bold': True, 'color': GREEN})],
    72, DARK, True, first=True)
sub = tbox(s, Inches(0.62), Inches(4.05), Inches(9.2), Inches(1.0))
par(sub, [("Gestion intelligente de cabinet médical — aide au diagnostic par IA, "
           "téléconsultation et automatisation du parcours de soins.",
           {'size': 18, 'color': GREY})], 18, GREY, first=True, ls=1.15)
hline(s, Inches(0.6), Inches(5.35), Inches(12.13), color=LINE, weight=1.3)
# presenters
cols = [("PRÉSENTÉ PAR", "Abdou SALOU ABDOU", 0.6, 3.0),
        ("PRÉSENTÉ PAR", "Kamara MACIRE", 3.7, 2.6),
        ("PRÉSENTÉ PAR", "Nouridine SAWADOGO", 6.4, 2.9),
        ("ENCADRÉ PAR", "Dr. Soumia CHOKRI", 9.5, 3.2)]
for lab, nm, x, w in cols:
    tf = tbox(s, Inches(x), Inches(5.6), Inches(w), Inches(0.8))
    par(tf, [(lab, {'size': 10, 'bold': True, 'color': GREEN})], 10, GREEN, True,
        first=True, sa=3)
    par(tf, [(nm, {'size': 14, 'bold': True, 'color': DARK})], 14, DARK, True)
# ISGA + date bottom
if os.path.exists(P("00-logo-isga.png")):
    s.shapes.add_picture(P("00-logo-isga.png"), Inches(0.6), Inches(6.75),
                         height=Inches(0.4))
dt = tbox(s, Inches(9.0), Inches(6.7), Inches(3.73), Inches(0.7))
par(dt, [("DATE", {'size': 10, 'bold': True, 'color': LGREY})], 10, LGREY, True,
    align=PP_ALIGN.RIGHT, first=True, sa=2)
par(dt, [("Juin 2026 · Casablanca, Maroc", {'size': 13, 'bold': True, 'color': DARK})],
    13, DARK, True, align=PP_ALIGN.RIGHT)

# =================== SLIDE 2 — CONTEXTE ===================
s = slide(); set_bg(s, BG)
eyebrow(s, "Contexte · Problématique"); title(s, [("Le constat", {})])
keyline(s, "La gestion d'un cabinet médical reste freinée par des outils éparpillés "
           "et l'absence d'aide à la décision.")
bullets(s, [
    ("Une gestion manuelle et dispersée :", "papier, fichiers et outils multiples."),
    ("Le temps administratif", "empiète sur le temps médical."),
    ("Risques d'erreurs et d'oublis :", "diagnostics, suivis, rendez-vous."),
    ("Aucun outil d'aide à la décision", "diagnostique au quotidien."),
    ("Communication patient limitée :", "rappels, transmission des documents."),
], w=6.0)
card(s, 7.0, 2.45, 5.73, 4.2, fill=WHITE, line=LINE)
ct = tbox(s, Inches(7.35), Inches(2.75), Inches(5.1), Inches(3.6))
par(ct, [("LE QUOTIDIEN FRAGMENTÉ", {'size': 12, 'bold': True, 'color': GREY})],
    12, GREY, True, first=True, sa=10)
for it in ["Papier", "Fichiers Excel", "Téléphone", "Agenda papier", "Notes éparses"]:
    par(ct, [("•  ", {'color': GREENL, 'bold': True, 'size': 15}),
             (it, {'size': 15, 'color': DARK})], 15, DARK, sa=7)
par(ct, [("0", {'size': 40, 'bold': True, 'color': GREEN})], 40, GREEN, True, sa=0)
par(ct, [("outil d'aide à la décision relié au dossier patient",
          {'size': 13, 'color': GREY})], 13, GREY)
footer(s, 2)

# =================== SLIDE 3 — OBJECTIFS ===================
s = slide(); set_bg(s, BG)
eyebrow(s, "Objectifs · Périmètre"); title(s, [("Notre réponse", {})])
keyline(s, "Une seule plateforme pour réunir gestion, diagnostic assisté, "
           "sécurité et communication.")
objs = [("Centraliser", "Toute la gestion du cabinet sur une plateforme web unique."),
        ("Assister", "Aider le médecin au diagnostic grâce à l'intelligence artificielle."),
        ("Sécuriser", "Accès aux données par rôles : admin, secrétaire, médecin, patient."),
        ("Digitaliser", "Ordonnances et comptes-rendus en PDF générés automatiquement."),
        ("Faciliter l'accès", "Téléconsultation et espace patient en ligne."),
        ("Automatiser", "Rappels et envois via WhatsApp et e-mail.")]
cw, ch, gx, gy = 3.84, 1.9, 0.3, 0.28
x0, y0 = 0.6, 2.5
for i, (lab, desc) in enumerate(objs):
    r, c = divmod(i, 3)
    x = x0 + c * (cw + gx); y = y0 + r * (ch + gy)
    card(s, x, y, cw, ch)
    tf = tbox(s, Inches(x + 0.25), Inches(y + 0.22), Inches(cw - 0.5), Inches(ch - 0.4))
    par(tf, [(lab, {'size': 17, 'bold': True, 'color': GREEN})], 17, GREEN, True,
        first=True, sa=6)
    par(tf, [(desc, {'size': 13, 'color': GREY})], 13, GREY, ls=1.1)
footer(s, 3)

# =================== SLIDE 4 — ARCHITECTURE ===================
std(4, "Conception · Technique", [("Architecture technique", {})],
    "Une architecture moderne, en microservices, entièrement conteneurisée.",
    [("Frontend", "React + Vite + Tailwind CSS."),
     ("API", "Django + Django REST Framework (REST sécurisée par JWT)."),
     ("Microservice IA", "Flask + Scikit-learn."),
     ("Base de données", "PostgreSQL."),
     ("Orchestration", "Docker / Docker Compose."),
     ("Intégrations", "n8n, Twilio/WhatsApp, Jitsi, Groq (LLM).")],
    images=[P("04-architecture-globale.png")], ib=(6.4, 2.3, 6.33, 4.45))

# =================== SLIDE 5 — CAS D'UTILISATION ===================
std(5, "Conception · Acteurs", [("Cas d'utilisation & rôles", {})],
    "Quatre rôles, des droits strictement cloisonnés.",
    [("Administrateur", "configuration, utilisateurs, audit."),
     ("Secrétaire", "patients, rendez-vous, planning."),
     ("Médecin", "consultations, ordonnances, assistance IA."),
     ("Patient", "espace personnel, rendez-vous, documents.")],
    images=[P("05-cas-utilisation.png")], ib=(6.4, 2.3, 6.33, 4.45))

# =================== SLIDE 6 — DIAGRAMME DE CLASSES ===================
std(6, "Conception · Données", [("Modèle de données", {})],
    "Un modèle de données clair et relationnel.",
    [("Entités clés :", "Utilisateur, Patient, Rendez-vous, Consultation, Ordonnance."),
     ("Un patient", "a plusieurs rendez-vous et consultations."),
     ("Une consultation", "génère une ordonnance et un compte-rendu.")],
    images=[P("06-diagramme-classe.png")], ib=(6.0, 2.25, 6.73, 4.55))

# =================== SLIDE 7 — SÉCURITÉ ===================
std(7, "Sécurité · Accès", [("Sécurité & contrôle d'accès", {})],
    "Chaque rôle ne voit que ce qui le concerne.",
    [("Authentification par JWT", "(SimpleJWT)."),
     ("Permissions par rôle", "au niveau de l'API (DRF)."),
     ("Cloisonnement strict", "des données ; tout accès non autorisé est bloqué."),
     ("Corbeille et traçabilité", "des actions.")],
    images=[P("07-controle-acces-role.png")])

# =================== SLIDE 8 — SECRÉTAIRE ===================
std(8, "Rôle · Secrétaire", [("Gestion patients & rendez-vous", {})],
    "Le quotidien de l'accueil, centralisé et fluide.",
    [("Dossiers patients", "complets, recherche rapide."),
     ("Création et suivi", "des rendez-vous."),
     ("Planning", "par jour et par médecin."),
     ("Statuts & corbeille", "de récupération.")],
    images=[P("08a-secretaire-patients.png"), P("08b-secretaire-planning-rdv.png")])

# =================== SLIDE 9 — MÉDECIN ===================
std(9, "Rôle · Médecin", [("Consultations", {})],
    "Une fiche patient complète, au service de la décision clinique.",
    [("Fiche patient 360°", ": informations, antécédents, historique."),
     ("Saisie structurée", "de la consultation."),
     ("Lien direct", "vers l'ordonnance et le compte-rendu."),
     ("Suivi de l'évolution", "du patient dans le temps.")],
    images=[P("09-medecin-fiche-patient.png")])

# =================== SLIDE 10 — IA (HERO SOMBRE) ===================
s = slide(); set_bg(s, DARKBG)
eyebrow(s, "Intelligence artificielle", color=VIOLET)
title(s, [("Assistance au diagnostic par ", {'color': WHITE}),
          ("IA", {'color': VIOLET})], color=WHITE)
keyline(s, "Le cœur de CuraMedical : aider le médecin à décider.", color=LGREY)
bullets(s, [
    ("Le médecin", "sélectionne les symptômes observés."),
    ("L'IA propose instantanément", "les maladies les plus probables."),
    ("Un appui à la décision", "— le médecin garde toujours le dernier mot.")],
    mark=VIOLET, head=WHITE, body=LGREY, y=2.6, sa=14)
img_fit(s, P("10-ia-diagnostic-suggestion.png"),
        Inches(6.55), Inches(2.35), Inches(6.18), Inches(4.35))
footer(s, 10, dark=True)

# =================== SLIDE 11 — PERFORMANCES IA ===================
s = slide(); set_bg(s, BG)
eyebrow(s, "Intelligence artificielle · Évaluation", color=VIOLET)
title(s, [("Performances du modèle IA", {})])
big = tbox(s, Inches(0.62), Inches(2.0), Inches(5.0), Inches(2.0))
par(big, [("77,3", {'size': 96, 'bold': True, 'color': VIOLET}),
          (" %", {'size': 44, 'bold': True, 'color': VIOLET})], 96, VIOLET, True, first=True)
sub = tbox(s, Inches(0.66), Inches(3.85), Inches(5.0), Inches(0.6))
par(sub, [("de précision sur le jeu de test", {'size': 17, 'color': GREY})],
    17, GREY, first=True)
stats = [("15 000", "cas cliniques (entraînement)"),
         ("377", "symptômes en entrée"),
         ("655", "maladies possibles"),
         ("Random Forest", "Scikit-learn · 30 arbres")]
sx, sy, sw, sh, g = 6.4, 2.0, 3.05, 1.55, 0.25
for i, (num, lab) in enumerate(stats):
    r, c = divmod(i, 2)
    x = sx + c * (sw + g); y = sy + r * (sh + g)
    card(s, x, y, sw, sh)
    tf = tbox(s, Inches(x + 0.25), Inches(y + 0.22), Inches(sw - 0.5), Inches(sh - 0.4))
    par(tf, [(num, {'size': 26, 'bold': True, 'color': DARK})], 26, DARK, True,
        first=True, sa=3)
    par(tf, [(lab, {'size': 12.5, 'color': GREY})], 12.5, GREY, ls=1.05)
band = card(s, 0.6, 5.55, 12.13, 1.0, fill=VIOLBG, line=VIOLBG)
bt = tbox(s, Inches(0.95), Inches(5.72), Inches(11.5), Inches(0.7),
          anchor=MSO_ANCHOR.MIDDLE)
par(bt, [("Mise en perspective : ", {'size': 14, 'bold': True, 'color': VIOLET}),
         ("655 classes possibles (baseline aléatoire ≈ 0,15 %). "
          "Une aide à la décision — jamais un diagnostic automatique.",
          {'size': 14, 'color': DARK})], 14, DARK, first=True, ls=1.1)
footer(s, 11)

# =================== SLIDE 12 — ORDONNANCES ===================
std(12, "Documents médicaux", [("Ordonnances & documents", {})],
    "De la saisie au document professionnel, en un clic.",
    [("Saisie structurée", "des médicaments et posologies."),
     ("Ordonnance PDF", "générée automatiquement."),
     ("Cachet et signature", "du médecin intégrés."),
     ("Comptes-rendus", "également exportés en PDF.")],
    images=[P("12a-saisie-ordonnance.png"), P("12b-ordonnance-pdf-genere.png")])

# =================== SLIDE 13 — TÉLÉCONSULTATION ===================
std(13, "Services avancés", [("Téléconsultation", {})],
    "L'accès aux soins à distance, en visioconférence sécurisée.",
    [("Visioconférence intégrée", "(Jitsi)."),
     ("Salle dédiée", "par rendez-vous."),
     ("Prise de notes en direct", "par le médecin."),
     ("Accès aux soins à distance", "pour le patient.")],
    images=[P("13-teleconsultation.png")])

# =================== SLIDE 14 — CHATBOT ===================
std(14, "Services avancés · IA", [("Chatbot IA", {})],
    "Interroger le cabinet en langage naturel.",
    [("Questions en français courant", ": statistiques, patients, rendez-vous."),
     ("Réponses instantanées", ": nombre de patients, RDV du jour…"),
     ("Propulsé par un LLM", "(Groq).")],
    images=[P("14-chatbot-ia.png")])

# =================== SLIDE 15 — WHATSAPP ===================
std(15, "Automatisation · n8n", [("Notifications WhatsApp", {})],
    "Rejoindre le patient là où il est, automatiquement.",
    [("Rappels de rendez-vous", "automatiques."),
     ("Envoi des ordonnances", "directement au patient."),
     ("Orchestration", "via workflow n8n + Twilio."),
     ("Moins d'absences", ", meilleure relation patient.")],
    images=[P("15-notifications-whatsapp.jpg")], ib=(7.4, 2.2, 4.6, 4.6))

# =================== SLIDE 16 — ESPACE PATIENT ===================
std(16, "Rôle · Patient", [("Espace patient", {})],
    "Le patient, acteur de son propre parcours de soins.",
    [("Inscription en ligne", "du patient."),
     ("Consultation de ses rendez-vous", "."),
     ("Accès à ses consultations", "et ordonnances."),
     ("Gestion de son profil", "personnel.")],
    images=[P("16-patient-espace-personnel.png")])

# =================== SLIDE 17 — CONCLUSION ===================
s = slide(); set_bg(s, BG)
eyebrow(s, "Conclusion"); title(s, [("Bilan & perspectives", {})])
keyline(s, "Un projet mené de bout en bout, de l'analyse au déploiement.")
card(s, 0.6, 2.5, 5.9, 4.0)
bt = tbox(s, Inches(0.95), Inches(2.8), Inches(5.3), Inches(3.5))
par(bt, [("BILAN", {'size': 14, 'bold': True, 'color': GREEN})], 14, GREEN, True,
    first=True, sa=12)
for b, n in [("Une plateforme", "fonctionnelle de bout en bout."),
             ("4 rôles", "et une IA opérationnelle."),
             ("Documents automatisés", "et intégrations réelles.")]:
    par(bt, [("✓  ", {'color': GREEN, 'bold': True, 'size': 15}),
             (b + " ", {'bold': True, 'color': DARK, 'size': 15}),
             (n, {'color': GREY, 'size': 15})], 15, DARK, sa=10, ls=1.1)
card(s, 6.83, 2.5, 5.9, 4.0)
pt = tbox(s, Inches(7.18), Inches(2.8), Inches(5.3), Inches(3.5))
par(pt, [("PERSPECTIVES", {'size': 14, 'bold': True, 'color': VIOLET})], 14, VIOLET,
    True, first=True, sa=12)
for b, n in [("Déploiement", "en production."),
             ("Amélioration continue", "du modèle."),
             ("Application mobile", "pour les patients."),
             ("Interopérabilité", "avec d'autres systèmes de santé.")]:
    par(pt, [("→  ", {'color': VIOLET, 'bold': True, 'size': 15}),
             (b + " ", {'bold': True, 'color': DARK, 'size': 15}),
             (n, {'color': GREY, 'size': 15})], 15, DARK, sa=10, ls=1.1)
footer(s, 17)

prs.save("CuraMedical-CC.pptx")
print("OK -> CuraMedical-CC.pptx :", len(prs.slides), "slides")
