# -*- coding: utf-8 -*-
"""Édite CuraMedical-CC.pptx EN PLACE (préserve les modifs visuelles de l'utilisateur) :
1) corrige la répétition de « PRÉSENTÉ PAR » sur la page de garde,
2) ajoute une slide « Plan » en position 2 (footer 01/17),
3) ajoute un indicateur de partie centré en pied de chaque slide numérotée.
La numérotation existante (02..17) reste inchangée -> la garde n'est pas comptée."""
import os, re, shutil
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE

BG=RGBColor(0xF7,0xF8,0xFA); WHITE=RGBColor(0xFF,0xFF,0xFF)
GREEN=RGBColor(0x2A,0x9B,0x69); DARK=RGBColor(0x0F,0x17,0x2A)
GREY=RGBColor(0x64,0x74,0x8B); LINE=RGBColor(0xE2,0xE8,0xF0)
LGREY=RGBColor(0x94,0xA3,0xB8); FONT="Segoe UI"
TOTAL=17
FN="CuraMedical-CC.pptx"

shutil.copyfile(FN, "CuraMedical-CC.bak.pptx")
prs=Presentation(FN)
BLANK=prs.slide_layouts[6]

def set_bg(s,rgb):
    f=s.background.fill; f.solid(); f.fore_color.rgb=rgb
def rect(s,x,y,w,h,fill=None,line=None,lw=Pt(1),rounded=False):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,x,y,w,h)
    if fill is None: sh.fill.background()
    else: sh.fill.solid(); sh.fill.fore_color.rgb=fill
    if line is None: sh.line.fill.background()
    else: sh.line.color.rgb=line; sh.line.width=lw
    sh.shadow.inherit=False; return sh
def hline(s,x,y,w,color=LINE,weight=1.0):
    return rect(s,x,y,w,Pt(weight),fill=color)
def tbox(s,x,y,w,h,anchor=MSO_ANCHOR.TOP):
    tb=s.shapes.add_textbox(x,y,w,h); tf=tb.text_frame
    tf.word_wrap=True; tf.auto_size=MSO_AUTO_SIZE.NONE; tf.vertical_anchor=anchor
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0; return tf
def par(tf,runs,size,color,bold=False,align=PP_ALIGN.LEFT,first=False,sa=None,ls=None):
    p=tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment=align
    if sa is not None: p.space_after=Pt(sa)
    if ls is not None: p.line_spacing=ls
    for t,o in runs:
        r=p.add_run(); r.text=t; f=r.font
        f.size=Pt(o.get('size',size)); f.bold=o.get('bold',bold)
        f.name=FONT; f.color.rgb=o.get('color',color)
    return p
def card(s,x,y,w,h,fill=WHITE,line=LINE):
    return rect(s,Inches(x),Inches(y),Inches(w),Inches(h),fill=fill,line=line,lw=Pt(1),rounded=True)
def footer(s,n):
    hline(s,Inches(0.6),Inches(7.02),Inches(12.13),color=LINE,weight=1)
    if os.path.exists("images/00-icone-curamedical.png"):
        s.shapes.add_picture("images/00-icone-curamedical.png",Inches(0.6),Inches(7.12),height=Inches(0.24))
    tf=tbox(s,Inches(0.95),Inches(7.13),Inches(3),Inches(0.3))
    par(tf,[("CuraMedical",{'size':10.5,'bold':True,'color':GREEN})],10.5,GREEN,True,first=True)
    t2=tbox(s,Inches(8.0),Inches(7.13),Inches(4.73),Inches(0.3))
    par(t2,[("%02d / %d"%(n,TOTAL),{'size':10.5,'color':GREY})],10.5,GREY,align=PP_ALIGN.RIGHT,first=True)

# ---------- 1) FIX COVER « PRÉSENTÉ PAR » ----------
cover=prs.slides[0]
labs=[sh for sh in cover.shapes if sh.has_text_frame and sh.text_frame.text.strip().startswith("PRÉSENTÉ PAR")]
labs.sort(key=lambda s:s.left)
for sh in labs[1:]:
    for p in sh.text_frame.paragraphs:
        if "PRÉSENTÉ PAR" in "".join(r.text for r in p.runs):
            for r in p.runs: r.text=""
            break
print("Cover: labels 'PRÉSENTÉ PAR' en double vidés:", len(labs)-1)

# ---------- 2) PLAN SLIDE ----------
plan=prs.slides.add_slide(BLANK); set_bg(plan,BG)
tf=tbox(plan,Inches(0.62),Inches(0.52),Inches(12),Inches(0.3))
par(tf,[("SOMMAIRE",{'size':12.5,'bold':True,'color':GREEN})],12.5,GREEN,True,first=True)
tt=tbox(plan,Inches(0.6),Inches(0.84),Inches(12.1),Inches(0.95))
par(tt,[("Plan de la présentation",{'size':38,'bold':True,'color':DARK})],38,DARK,True,first=True)
kt=tbox(plan,Inches(0.62),Inches(1.66),Inches(11.8),Inches(0.6))
par(kt,[("Six parties, du besoin métier jusqu'à l'intelligence artificielle.",{'size':17,'color':GREY})],17,GREY,first=True)
items=[("01","Contexte & objectifs","Le constat et les objectifs du projet."),
 ("02","Architecture & conception","Stack technique, cas d'usage, modèle de données."),
 ("03","L'application par rôle","Sécurité, secrétaire, médecin."),
 ("04","L'intelligence artificielle","Aide au diagnostic et performances du modèle."),
 ("05","Services & automatisation","Ordonnances, téléconsultation, chatbot, WhatsApp."),
 ("06","Espace patient & conclusion","Portail patient, bilan et perspectives.")]
cw,ch,gx,gy=5.9,1.3,0.33,0.2; x0,y0=0.6,2.5
for i,(num,t,d) in enumerate(items):
    r,c=divmod(i,2); x=x0+c*(cw+gx); y=y0+r*(ch+gy)
    card(plan,x,y,cw,ch)
    nb=tbox(plan,Inches(x+0.3),Inches(y),Inches(1.1),Inches(ch),anchor=MSO_ANCHOR.MIDDLE)
    par(nb,[(num,{'size':30,'bold':True,'color':GREEN})],30,GREEN,True,first=True)
    bt=tbox(plan,Inches(x+1.45),Inches(y+0.24),Inches(cw-1.7),Inches(ch-0.45))
    par(bt,[(t,{'size':16,'bold':True,'color':DARK})],16,DARK,True,first=True,sa=3)
    par(bt,[(d,{'size':11.5,'color':GREY})],11.5,GREY,ls=1.03)
footer(plan,1)
# déplacer le plan en position 2 (index 1)
lst=prs.slides._sldIdLst; ids=list(lst)
lst.remove(ids[-1]); lst.insert(1,ids[-1])
print("Plan ajouté en position 2.")

# ---------- 3) INDICATEUR DE PARTIE ----------
SEC={1:"SOMMAIRE"}
for n in (2,3): SEC[n]="PARTIE 1 · CONTEXTE & OBJECTIFS"
for n in (4,5,6): SEC[n]="PARTIE 2 · ARCHITECTURE & CONCEPTION"
for n in (7,8,9): SEC[n]="PARTIE 3 · L'APPLICATION PAR RÔLE"
for n in (10,11): SEC[n]="PARTIE 4 · INTELLIGENCE ARTIFICIELLE"
for n in (12,13,14,15): SEC[n]="PARTIE 5 · SERVICES & AUTOMATISATION"
for n in (16,17): SEC[n]="PARTIE 6 · ESPACE PATIENT & CONCLUSION"

def slide_num(slide):
    for sh in slide.shapes:
        if sh.has_text_frame:
            m=re.match(r'^(\d+)\s*/\s*\d+$', sh.text_frame.text.strip())
            if m: return int(m.group(1))
    return None

added=0
for slide in prs.slides:
    n=slide_num(slide)
    if n is None or n not in SEC: continue
    col=LGREY if n==10 else GREY   # slide 10 = fond sombre
    tf=tbox(slide,Inches(4.0),Inches(7.14),Inches(5.33),Inches(0.3),anchor=MSO_ANCHOR.TOP)
    par(tf,[(SEC[n],{'size':9,'bold':True,'color':col})],9,col,True,align=PP_ALIGN.CENTER,first=True)
    added+=1
print("Indicateurs de partie ajoutés:",added)

prs.save(FN)
print("OK -> %s : %d slides (sauvegarde: CuraMedical-CC.bak.pptx)"%(FN,len(prs.slides)))
