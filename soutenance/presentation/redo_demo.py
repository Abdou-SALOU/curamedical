# -*- coding: utf-8 -*-
"""(1) Retire liens/boutons/transitions/animations ajoutes precedemment.
   (2) Refait la diapo "Demonstration" en maquette navigateur + vraie capture."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from PIL import Image

DOCX = "CuraMedical-Présentation.pptx"
SHOT = "../screenshots/curamedical-rapport/16-medecin-nouvelle-consultation-ia.png"

BG=RGBColor(0xF7,0xF8,0xFA); WHITE=RGBColor(0xFF,0xFF,0xFF); GREEN=RGBColor(0x2A,0x9B,0x69)
DARK=RGBColor(0x0F,0x17,0x2A); GREY=RGBColor(0x64,0x74,0x8B); LINE=RGBColor(0xE2,0xE8,0xF0)
CHROME=RGBColor(0xED,0xF1,0xF5); CHROMEB=RGBColor(0xD7,0xDE,0xE6); SHADOW=RGBColor(0xDA,0xE0,0xE8)
FONT="Segoe UI"; TOTAL=18

prs = Presentation(DOCX)

# ---------- (1) NETTOYAGE interactivite ----------
nhl=nt=ntm=npb=0
for s in prs.slides:
    sld=s._element
    for hl in sld.findall('.//'+qn('a:hlinkClick')): hl.getparent().remove(hl); nhl+=1
    for t in sld.findall(qn('p:transition')): t.getparent().remove(t); nt+=1
    for tm in sld.findall(qn('p:timing')): tm.getparent().remove(tm); ntm+=1
for s in prs.slides:
    for sh in list(s.shapes):
        if '__PLAN__' in (sh.name or ''):
            sh._element.getparent().remove(sh._element); npb+=1
print(f"Nettoyage -> hyperliens:{nhl} transitions:{nt} animations:{ntm} boutonsPlan:{npb}")

# ---------- helpers ----------
def set_bg(s,rgb):
    f=s.background.fill; f.solid(); f.fore_color.rgb=rgb
def rect(s,x,y,w,h,fill=None,line=None,lw=Pt(1),rounded=False):
    shp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,x,y,w,h)
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb=fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb=line; shp.line.width=lw
    shp.shadow.inherit=False; return shp
def oval(s,x,y,d,fill):
    shp=s.shapes.add_shape(MSO_SHAPE.OVAL,x,y,d,d)
    shp.fill.solid(); shp.fill.fore_color.rgb=fill; shp.line.fill.background()
    shp.shadow.inherit=False; return shp
def hline(s,x,y,w,color=LINE,weight=1.0): return rect(s,x,y,w,Pt(weight),fill=color)
def tbox(s,x,y,w,h,anchor=MSO_ANCHOR.TOP,wrap=True):
    tb=s.shapes.add_textbox(x,y,w,h); tf=tb.text_frame
    tf.word_wrap=wrap; tf.auto_size=MSO_AUTO_SIZE.NONE; tf.vertical_anchor=anchor
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0; return tf
def par(tf,runs,size,color,bold=False,align=PP_ALIGN.LEFT,first=False,sa=None,ls=None):
    p=tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment=align
    if sa is not None: p.space_after=Pt(sa)
    if ls is not None: p.line_spacing=ls
    if isinstance(runs,str): runs=[(runs,{})]
    for t,o in runs:
        r=p.add_run(); r.text=t; f=r.font
        f.size=Pt(o.get('size',size)); f.bold=o.get('bold',bold)
        f.name=FONT; f.color.rgb=o.get('color',color)
    return p
def footer(s,n):
    hline(s,Inches(0.6),Inches(7.02),Inches(12.13),color=LINE,weight=1)
    if os.path.exists("images/00-icone-curamedical.png"):
        s.shapes.add_picture("images/00-icone-curamedical.png",Inches(0.6),Inches(7.12),height=Inches(0.24))
    par(tbox(s,Inches(0.95),Inches(7.13),Inches(5),Inches(0.3)),
        [("CuraMedical",{'size':10.5,'bold':True,'color':GREEN})],10.5,GREEN,True,first=True)
    par(tbox(s,Inches(8.0),Inches(7.13),Inches(4.73),Inches(0.3)),
        [("%02d / %d"%(n,TOTAL),{'size':10.5,'color':GREY})],10.5,GREY,align=PP_ALIGN.RIGHT,first=True)
    par(tbox(s,Inches(5.0),Inches(7.13),Inches(3.2),Inches(0.3),anchor=MSO_ANCHOR.MIDDLE),
        [("PARTIE 6 · ESPACE PATIENT & CONCLUSION",{'size':9,'bold':True,'color':GREY})],9,GREY,True,align=PP_ALIGN.CENTER,first=True)

# ---------- (2) trouver la diapo demo et la vider ----------
demo=None
for s in prs.slides:
    for sh in s.shapes:
        if sh.has_text_frame and "DÉMONSTRATION" in sh.text_frame.text.upper() and "DÉMO" not in sh.text_frame.text.upper():
            demo=s; break
    if demo is None:
        for sh in s.shapes:
            if sh.has_text_frame and "Place à la démonstration" in sh.text_frame.text:
                demo=s; break
    if demo: break
assert demo is not None, "diapo demo introuvable"
for sh in list(demo.shapes):
    sh._element.getparent().remove(sh._element)
set_bg(demo,BG)

# entete
par(tbox(demo,Inches(0.62),Inches(0.52),Inches(12),Inches(0.3)),
    [("DÉMONSTRATION · LIVE",{'size':12.5,'bold':True,'color':GREEN})],12.5,GREEN,True,first=True)
par(tbox(demo,Inches(0.6),Inches(0.84),Inches(12.1),Inches(0.95)),
    [("Place à la démonstration",{})],38,DARK,True,first=True)
par(tbox(demo,Inches(0.62),Inches(1.66),Inches(11.8),Inches(0.6)),
    [("L'application en conditions réelles, du rendez-vous à l'ordonnance.",{'size':17,'color':GREY})],17,GREY,first=True)

# puces scenario (gauche)
items=[
    ("Rendez-vous et planning","(secrétaire)."),
    ("Consultation assistée par l'IA","(symptômes, top diagnostics)."),
    ("Ordonnance PDF","(cachet et signature)."),
    ("Notification WhatsApp","envoyée au patient."),
    ("Téléconsultation et espace patient","en ligne."),
]
tf=tbox(demo,Inches(0.62),Inches(2.55),Inches(4.85),Inches(4.1))
fst=True
for b,n in items:
    par(tf,[("›  ",{'color':GREEN,'bold':True,'size':15.5}),
            (b+" ",{'bold':True,'color':DARK,'size':15.5}),
            (n,{'color':GREY,'size':15.5})],15.5,DARK,first=fst,sa=15,ls=1.12); fst=False

# ---------- maquette navigateur (droite) ----------
WX, WY, WW = 5.62, 2.12, 7.4          # fenetre
pad = 0.12
shot_w = WW - 2*pad
shot_h = shot_w / (16/9)              # captures 16:9
chrome_h = 0.42
WH = chrome_h + shot_h + pad
# ombre douce
rect(demo,Inches(WX+0.07),Inches(WY+0.09),Inches(WW),Inches(WH),fill=SHADOW,rounded=True)
# fenetre
rect(demo,Inches(WX),Inches(WY),Inches(WW),Inches(WH),fill=CHROME,line=CHROMEB,lw=Pt(1),rounded=True)
# pastilles
for k,col in enumerate([RGBColor(0xFF,0x5F,0x57),RGBColor(0xFE,0xBC,0x2E),RGBColor(0x28,0xC8,0x40)]):
    oval(demo,Inches(WX+0.22+k*0.22),Inches(WY+0.15),Inches(0.12),col)
# barre d'URL
rect(demo,Inches(WX+0.95),Inches(WY+0.09),Inches(WW-1.55),Inches(0.25),fill=WHITE,line=CHROMEB,lw=Pt(0.75),rounded=True)
par(tbox(demo,Inches(WX+1.12),Inches(WY+0.10),Inches(WW-1.9),Inches(0.23),anchor=MSO_ANCHOR.MIDDLE),
    [("localhost:5173",{'size':9.5,'bold':True,'color':GREEN}),
     ("/consultations",{'size':9.5,'color':GREY})],9.5,GREY,first=True)
# capture reelle
if os.path.exists(SHOT):
    pic=demo.shapes.add_picture(SHOT,Inches(WX+pad),Inches(WY+chrome_h),Inches(shot_w),Inches(shot_h))
    pic.line.color.rgb=CHROMEB; pic.line.width=Pt(0.75); pic.shadow.inherit=False
else:
    print("!! capture absente:",SHOT)

footer(demo,17)

prs.save(DOCX)
print("SAVED. Fenetre maquette: %.2f x %.2f in (capture %.2f x %.2f)"%(WW,WH,shot_w,shot_h))
